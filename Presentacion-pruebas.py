from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]


title.text = "Hola Mundo"
subtitle.text = "python es mi héroe"

prs.save("Presentacion.pptx")